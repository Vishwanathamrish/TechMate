from flask import Flask, request, jsonify, send_file, render_template, redirect, url_for, flash
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from dotenv import load_dotenv
import mysql.connector
import bcrypt
import PyPDF2
import nltk
import re
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import base64

# Load environment variables from .env file
load_dotenv()
from qdrant_client import QdrantClient
from qdrant_client.http.models import PointStruct, VectorParams, Distance
import uuid
from datetime import datetime, timezone
import groq
import logging

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Initialize Groq client
import os

# Load API key from .env file
GROQ_API_KEY = os.getenv('GROQ_API_KEY', '').strip()
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY not found in .env file")
try:
    groq_client = groq.Groq(
        api_key=GROQ_API_KEY
    )
    # Test the API key with a simple request
    test_completion = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": "test"}],
        max_tokens=1
    )
    logging.info("GroQ API key validated successfully")
except groq.AuthenticationError as e:
    logging.error(f"GroQ API authentication error: {str(e)}")
    raise ValueError(f"Invalid or expired GroQ API key. Please check your API key.")
except Exception as e:
    logging.error(f"Error initializing GroQ client: {str(e)}")
    raise ValueError(f"Failed to initialize GroQ client: {str(e)}")

# Download required NLTK data
try:
    nltk.download('punkt', quiet=True)
    nltk.download('averaged_perceptron_tagger', quiet=True)
except Exception as e:
    logging.error(f"Error downloading NLTK data: {str(e)}")
    raise

app = Flask(__name__)
import secrets
# Load SECRET_KEY from environment (.env). If not set, generate a secure random key for dev.
# Warning: the generated fallback changes on each restart — set SECRET_KEY persistently for production.
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY') or secrets.token_hex(32)

# Initialize Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

# Configure MySQL
db_config = {
    'host': 'localhost',
    'user': 'root',  # Replace with your MySQL username
    'password': 'Vishwanath1604@',  # Replace with your MySQL password
    'database': 'professor_bot_db'
}

try:
    db = mysql.connector.connect(**db_config)
    cursor = db.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INT AUTO_INCREMENT PRIMARY KEY,
            username VARCHAR(50) UNIQUE NOT NULL,
            password_hash VARCHAR(255) NOT NULL
        )
    ''')
    db.commit()
except Exception as e:
    logging.error(f"Error setting up MySQL database: {str(e)}")
    raise

# User class for Flask-Login
class User(UserMixin):
    def __init__(self, id, username):
        self.id = id
        self.username = username

@login_manager.user_loader
def load_user(user_id):
    try:
        cursor = db.cursor()
        cursor.execute("SELECT id, username FROM users WHERE id = %s", (user_id,))
        user = cursor.fetchone()
        if user:
            return User(id=user[0], username=user[1])
        return None
    except Exception as e:
        logging.error(f"Error loading user: {str(e)}")
        return None

# No configuration needed for Groq client as it's initialized above

# Initialize Qdrant client (currently bypassed)
qdrant_client = QdrantClient(url="http://localhost:6333")

# Qdrant collection name
COLLECTION_NAME = "professor_documents"

# Global variables
study_plan = []
mcq_questions = []
syllabus = []
current_document_id = None
current_document_text = ""
current_ppt_data = None  # Store the generated PPT data

def setup_qdrant_collection():
    """Set up Qdrant collection if it doesn't exist."""
    try:
        collections = qdrant_client.get_collections()
        if COLLECTION_NAME not in [c.name for c in collections.collections]:
            qdrant_client.create_collection(
                collection_name=COLLECTION_NAME,
                vectors_config=VectorParams(size=1, distance=Distance.COSINE)
            )
    except Exception as e:
        logging.error(f"Error setting up Qdrant collection: {str(e)}")

def chunk_document(text, chunk_size=500):
    """Split document text into smaller chunks for better retrieval."""
    try:
        sentences = nltk.sent_tokenize(text)
        chunks = []
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) <= chunk_size:
                current_chunk += sentence + " "
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + " "
        if current_chunk:
            chunks.append(current_chunk.strip())
        return chunks
    except Exception as e:
        logging.error(f"Error chunking document: {str(e)}")
        return []

def store_document_in_qdrant(text, document_id):
    """Store document chunks in Qdrant with valid UUID point IDs."""
    try:
        chunks = chunk_document(text)
        points = []
        for idx, chunk in enumerate(chunks):
            point_id = str(uuid.uuid4())
            point = PointStruct(
                id=point_id,
                vector=[random.random()],
                payload={
                    "text": chunk,
                    "document_id": document_id,
                    "chunk_idx": idx,
                    "uploaded_at": datetime.now(timezone.utc).isoformat()
                }
            )
            points.append(point)
        qdrant_client.upsert(
            collection_name=COLLECTION_NAME,
            points=points
        )
    except Exception as e:
        return f"Error storing document in Qdrant: {str(e)}"
    return None

def retrieve_relevant_chunks(query, document_id, top_k=3):
    """Retrieve relevant chunks from Qdrant based on the query."""
    try:
        points = qdrant_client.scroll(
            collection_name=COLLECTION_NAME,
            scroll_filter={"must": [{"key": "document_id", "match": {"value": document_id}}]},
            limit=100
        )[0]
        if not points:
            return []

        query_words = set(query.lower().split())
        scored_chunks = []
        for point in points:
            chunk_text = point.payload["text"]
            chunk_words = set(chunk_text.lower().split())
            common_words = query_words.intersection(chunk_words)
            score = len(common_words) / len(query_words) if query_words else 0
            scored_chunks.append((chunk_text, score))

        scored_chunks.sort(key=lambda x: x[1], reverse=True)
        return [chunk for chunk, score in scored_chunks[:top_k] if score > 0]
    except Exception as e:
        return [f"Error retrieving chunks from Qdrant: {str(e)}"]

def get_document_from_qdrant(document_id):
    """Retrieve and reconstruct full document text from Qdrant."""
    try:
        points = qdrant_client.scroll(
            collection_name=COLLECTION_NAME,
            scroll_filter={"must": [{"key": "document_id", "match": {"value": document_id}}]},
            limit=100
        )[0]
        if not points:
            return None
        points.sort(key=lambda p: p.payload["chunk_idx"])
        document_text = " ".join(point.payload["text"] for point in points)
        return document_text
    except Exception as e:
        return f"Error retrieving document from Qdrant: {str(e)}"

def read_pdf(file_stream):
    """Read text from a PDF file stream."""
    try:
        pdf = PyPDF2.PdfReader(file_stream)
        document_text = ""
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                document_text += text + " "
        if not document_text.strip():
            return None, "Error: No text found in the PDF."
        return document_text, None
    except Exception as e:
        logging.error(f"Error reading PDF: {str(e)}")
        return None, f"Error reading PDF: {str(e)}"

def generate_study_plan(text):
    """Generate a study plan based strictly on document content."""
    global study_plan
    study_plan = []
    if not text:
        return ["Error: No document available."]
    
    try:
        sentences = nltk.sent_tokenize(text)
        if len(sentences) < 5:
            return ["Error: Document is too short to generate a study plan."]
        
        words = nltk.word_tokenize(text.lower())
        words = [w for w in words if w.isalpha()]
        freq = nltk.FreqDist(words)
        common_words = [word for word, count in freq.most_common(20) if len(word) > 5]
        
        for i, topic in enumerate(common_words[:5], 1):
            study_plan.append(f"Week {i}: Study '{topic}' - Focus on sections mentioning this term in the document.")
        return study_plan if study_plan else ["No key topics identified in the document."]
    except Exception as e:
        logging.error(f"Error generating study plan: {str(e)}")
        return [f"Error generating study plan: {str(e)}"]

def generate_syllabus_with_groq(text):
    """Generate a detailed syllabus using GroQ with llama-3.3-70b-versatile model."""
    global syllabus
    syllabus = []
    if not text:
        return ["Error: No document available to generate syllabus."]

    text = text[:5000]
    prompt = f"""Given the following document content, create a detailed syllabus for a student course. Include a course overview, 3–5 overall learning objectives, and 4–6 main topics with subtopics, timelines, prerequisites, estimated study hours, key resources from the document, learning outcomes, and assessment ideas. Focus on key concepts, organize logically, and ensure content is derived strictly from the document.

Content:
{text}

Output format:
Course Overview: [Summary of the course in 100–150 characters.]

Learning Objectives:
- [Objective 1, under 100 characters]
- [Objective 2, under 100 characters]
- [Objective 3, under 100 characters]

Syllabus Outline:
- Topic 1: [Main Topic]
  - Subtopic 1.1: [Subtopic, under 80 characters]
  - Subtopic 1.2: [Subtopic, under 80 characters]
  - Timeline: [e.g., Week 1–2]
  - Prerequisites: [e.g., Basic biology knowledge, or None, under 80 characters]
  - Estimated Study Hours: [e.g., 5 hours]
  - Key Resources: [Document sections or terms to focus on, under 120 characters]
  - Learning Outcome: [Outcome for this topic, under 100 characters]
  - Assessment: [Assessment idea, under 80 characters]
- Topic 2: [Main Topic]
  - Subtopic 2.1: [Subtopic, under 80 characters]
  - Subtopic 2.2: [Subtopic, under 80 characters]
  - Timeline: [e.g., Week 3]
  - Prerequisites: [e.g., Topic 1 knowledge, under 80 characters]
  - Estimated Study Hours: [e.g., 4 hours]
  - Key Resources: [Document sections or terms to focus on, under 120 characters]
  - Learning Outcome: [Outcome for this topic, under 100 characters]
  - Assessment: [Assessment idea, under 80 characters]
"""
    try:
        completion = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educational content creator specializing in creating detailed course syllabi."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        syllabus_text = completion.choices[0].message.content.strip()
        syllabus = syllabus_text.split("\n")
        syllabus = [line.strip() for line in syllabus if line.strip()]
        return syllabus if syllabus else ["No syllabus generated."]
    except Exception as e:
        logging.error(f"Error generating syllabus with GroQ: {str(e)}")
        return [f"Error generating syllabus with GroQ: {str(e)}"]

def generate_mcq(text):
    """Generate 10 MCQ questions strictly from document content without POS tagging."""
    global mcq_questions
    mcq_questions = []
    
    if not text or not isinstance(text, str):
        logging.error("Invalid or empty text provided for MCQ generation")
        return [{"question": "Error: No valid document content available.", "options": [], "correct": ""}]
    
    try:
        # Clean and preprocess text
        text = ' '.join(text.split())  # Normalize whitespace
        sentences = nltk.sent_tokenize(text)
        logging.info(f"Found {len(sentences)} sentences in document")
        
        # Filter key sentences - must be complete and meaningful
        key_sentences = [s for s in sentences if len(s.split()) >= 10 
                        and '?' not in s 
                        and len(s) <= 200  # Not too long
                        and any(c.isalpha() for c in s)]  # Contains letters
        
        if len(key_sentences) < 5:
            logging.warning(f"Insufficient key sentences found: {len(key_sentences)}")
            return [{"question": "Error: Document content is not suitable for MCQ generation.", 
                    "options": [], "correct": ""}]
        
        # Extract and filter key terms
        words = nltk.word_tokenize(text.lower())
        words = [w for w in words if w.isalpha() and len(w) > 3 
                and not w in ['this', 'that', 'these', 'those', 'there', 'where', 'when']]
        freq = nltk.FreqDist(words)
        key_terms = [word for word, count in freq.most_common(100) 
                    if count >= 2 and len(word) >= 4]  # Term must appear at least twice
        
        logging.info(f"Found {len(key_terms)} potential key terms")
        
        questions_generated = 0
        for sentence in key_sentences:
            if questions_generated >= 10:
                break
            
            sentence_lower = sentence.lower()
            # Find meaningful key terms in the sentence
            matching_terms = [term for term in key_terms 
                            if term in sentence_lower 
                            and not any(term in t for t in key_terms if t != term)]  # Not a substring of other terms
            
            if not matching_terms:
                continue
            
            # Use the most meaningful (longest) term
            key_term = max(matching_terms, key=len)
            
            # Create the question with proper context
            pattern = r'\b' + re.escape(key_term) + r'\b'
            if not re.search(pattern, sentence, re.IGNORECASE):
                continue
            
            question = re.sub(pattern, "______", sentence, count=1, flags=re.IGNORECASE)
            question = f"Question {questions_generated + 1}: {question}"
            
            # Generate plausible distractors
            # Filter to terms of similar length and type
            similar_terms = [t for t in key_terms 
                           if t != key_term 
                           and abs(len(t) - len(key_term)) <= 3  # Similar length
                           and not (t in key_term or key_term in t)]  # Not substring of each other
            
            if len(similar_terms) < 3:
                continue
            
            # Select distractors randomly but consistently
            random.seed(hash(sentence + key_term))  # Deterministic for same input
            distractors = random.sample(similar_terms, min(3, len(similar_terms)))
            options = [key_term] + distractors
            random.shuffle(options)
            
            mcq_questions.append({
                "question": question,
                "options": options,
                "correct": key_term
            })
            questions_generated += 1
            logging.debug(f"Generated MCQ {questions_generated}: {question}")
        
        if not mcq_questions:
            logging.warning("No MCQs could be generated")
            return [{"question": "Could not generate MCQs from the document content.", 
                    "options": [], "correct": ""}]
        
        logging.info(f"Successfully generated {len(mcq_questions)} MCQs")
        return mcq_questions
    except Exception as e:
        logging.error(f"Error generating MCQs: {str(e)}")
        return [{"question": f"Error generating MCQs: {str(e)}", "options": [], "correct": ""}]

def create_ppt_from_slides():
    """Create a PowerPoint presentation using content from the uploaded PDF."""
    try:
        global current_document_text
        if not current_document_text:
            return None, "Error: No document content available (please upload a PDF first)."

        sentences = nltk.sent_tokenize(current_document_text)
        if len(sentences) < 5:
            return None, "Error: Document is too short to generate a presentation."

        words = nltk.word_tokenize(current_document_text.lower())
        words = [w for w in words if w.isalpha() and len(w) > 5]
        freq = nltk.FreqDist(words)
        common_words = [word for word, count in freq.most_common(5)]

        if not common_words:
            return None, "Error: No key topics identified in the document."

        slides = []
        for topic in common_words:
            topic_sentences = [s for s in sentences if topic.lower() in s.lower()]
            if not topic_sentences:
                continue

            content = []
            for sentence in topic_sentences[:5]:
                bullet = sentence.strip()[:97] + "..." if len(sentence) > 97 else sentence.strip()
                content.append(bullet)
            if content:
                slides.append({
                    "title": f"Exploring {topic.capitalize()}",
                    "content": content
                })

        if not slides:
            return None, "Error: No suitable content found for slides."

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Presentation Based on Uploaded Document"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"

        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content_box = slide.placeholders[1]
            title.text = slide_data["title"]
            content_box.text = "\n".join(f"• {point}" for point in slide_data["content"])

        output = BytesIO()
        prs.save(output)
        return output.getvalue(), None
    except Exception as e:
        logging.error(f"Error creating PPT: {str(e)}")
        return None, f"Error creating PPT: {str(e)}"

def create_ppt_with_rag(topic, document_id):
    """Create a PowerPoint presentation using document content and GroQ."""
    global current_document_text
    if not current_document_text:
        return None, "Error: No document content available (Qdrant storage skipped)."

    context = current_document_text[:1000]
    prompt = f"Based on the following document content, generate content for a PowerPoint presentation about '{topic}'. Provide a title and 3 slide descriptions (each with a slide title and content, limited to 500 characters):\n\nContext:\n{context}\n\nOutput format:\nTitle: [Presentation Title]\nSlide 1: [Slide Title] - [Content]\nSlide 2: [Slide Title] - [Content]\nSlide 3: [Slide Title] - [Content]"

    try:
        completion = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at creating educational presentations."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=1000
        )
        ppt_content = completion.choices[0].message.content
    except Exception as e:
        return None, f"Error generating PPT content with Gemini: {str(e)}"

    lines = ppt_content.strip().split("\n")
    if len(lines) < 4:
        return None, "Error: Invalid PPT content generated by Gemini."

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    presentation_title = lines[0].replace("Title: ", "").strip()
    title.text = presentation_title
    subtitle.text = "Based on Uploaded Document"

    for i in range(1, 4):
        slide_info = lines[i].split(" - ", 1)
        if len(slide_info) != 2:
            continue
        slide_title, slide_content = slide_info
        slide_title = slide_title.replace(f"Slide {i}: ", "").strip()
        slide_content = slide_content.strip()

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content_box = slide.placeholders[1]
        title.text = slide_title
        content_box.text = slide_content[:500]

    output = BytesIO()
    prs.save(output)
    return output.getvalue(), None

def clarify_doubts_with_rag(question, document_id, response_type):
    """Clarify doubts using document content, syllabus, and GroQ."""
    if "i can't understand" in question.lower():
        return jsonify({"message": "Here is the updated code.", "status": "success"})

    global current_document_text
    if not current_document_text:
        return f"Error: No document content available (Qdrant storage skipped)."

    context = current_document_text[:1000]
    syllabus_context = "\n".join(syllabus) if syllabus else "No syllabus available."
    if response_type == "short":
        prompt = f"Based on the following document content and syllabus, provide a short, student-friendly explanation for the question: '{question}'. Keep the explanation under 100 characters.\n\nDocument Content:\n{context}\n\nSyllabus:\n{syllabus_context}\n\nOutput format:\nHere's a short explanation: [Explanation]."
    elif response_type == "detailed":
        prompt = f"Based on the following document content and syllabus, provide a detailed, student-friendly explanation for the question: '{question}'. Keep the explanation under 300 characters.\n\nDocument Content:\n{context}\n\nSyllabus:\n{syllabus_context}\n\nOutput format:\nHere's a detailed explanation: [Explanation]."
    else:
        prompt = f"Based on the following document content and syllabus, provide a simple explanation for the question: '{question}', and include an example. Keep the explanation under 150 characters and the example under 100 characters.\n\nDocument Content:\n{context}\n\nSyllabus:\n{syllabus_context}\n\nOutput format:\nHere's a simple explanation: [Explanation]. For example: [Example]."

    try:
        completion = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are a knowledgeable teacher explaining concepts to students in a clear and friendly way."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=500
        )
        return completion.choices[0].message.content
    except Exception as e:
        logging.error(f"Error generating explanation with Gemini: {str(e)}")
        return f"Error generating explanation with Gemini: {str(e)}"

@app.route('/login', methods=['GET', 'POST'])
def login():
    """Handle user login."""
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        try:
            username = request.form.get('username')
            password = request.form.get('password')
            if not username or not password:
                flash('Username and password are required.', 'error')
                return redirect(url_for('login'))

            cursor = db.cursor()
            cursor.execute("SELECT id, username, password_hash FROM users WHERE username = %s", (username,))
            user = cursor.fetchone()
            
            if user and bcrypt.checkpw(password.encode('utf-8'), user[2].encode('utf-8')):
                user_obj = User(id=user[0], username=user[1])
                login_user(user_obj)
                flash('Login successful!', 'success')
                return redirect(url_for('index'))
            else:
                flash('Invalid username or password.', 'error')
                return redirect(url_for('login'))
        except Exception as e:
            logging.error(f"Error during login: {str(e)}")
            flash('An error occurred during login.', 'error')
            return redirect(url_for('login'))
    
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    """Handle user registration."""
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        try:
            username = request.form.get('username')
            password = request.form.get('password')
            if not username or not password:
                flash('Username and password are required.', 'error')
                return redirect(url_for('register'))

            # Check if username already exists
            cursor = db.cursor()
            cursor.execute("SELECT id FROM users WHERE username = %s", (username,))
            if cursor.fetchone():
                flash('Username already exists.', 'error')
                return redirect(url_for('register'))

            # Hash password
            password_hash = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
            
            # Insert new user
            cursor.execute("INSERT INTO users (username, password_hash) VALUES (%s, %s)", (username, password_hash))
            db.commit()
            flash('Registration successful! Please log in.', 'success')
            return redirect(url_for('login'))
        except Exception as e:
            logging.error(f"Error during registration: {str(e)}")
            flash('An error occurred during registration.', 'error')
            return redirect(url_for('register'))
    
    return render_template('register.html')

@app.route('/logout')
@login_required
def logout():
    """Handle user logout."""
    logout_user()
    flash('Logged out successfully.', 'success')
    return redirect(url_for('login'))

@app.route('/')
@login_required
def index():
    """Serve the main chatbot interface."""
    try:
        setup_qdrant_collection()
        return render_template('index.html')
    except Exception as e:
        logging.error(f"Error rendering index.html: {str(e)}")
        return jsonify({"message": f"Error rendering UI: {str(e)}", "status": "error"}), 500

@app.route('/upload_pdf', methods=['POST'])
@login_required
def upload_pdf():
    """Handle PDF upload, store in Qdrant, generate study plan and syllabus."""
    try:
        logging.info("Received PDF upload request")
        if 'file' not in request.files:
            logging.warning("No file uploaded in request")
            return jsonify({"message": "No file uploaded", "status": "error"}), 400

        file = request.files['file']
        logging.info(f"Processing file: {file.filename}")
        
        text, error = read_pdf(file.stream)
        if error:
            logging.error(f"Failed to read PDF: {error}")
            return jsonify({"message": error, "status": "error"}), 400
        
        global current_document_id, current_document_text
        current_document_id = str(uuid.uuid4())
        current_document_text = text
        logging.info(f"PDF text extracted, document ID: {current_document_id}")
        
        plan = generate_study_plan(text)
        logging.info("Study plan generated")

        syllabus_result = generate_syllabus_with_groq(text)
        logging.info("Syllabus generation attempted")
        
        # Treat only items that start with 'Error' (case-insensitive) or 'Error:' as actual errors.
        import re
        def is_error_item(item):
            try:
                s = str(item).strip()
                return bool(re.match(r'^(?i:error)(:?\b|:)', s))
            except Exception:
                return False

        if isinstance(syllabus_result, list) and any(is_error_item(item) for item in syllabus_result):
            error_msg = next((item for item in syllabus_result if is_error_item(item)), "Unknown error")
            logging.error(f"Syllabus generation failed: {error_msg}")
            return jsonify({"message": str(error_msg), "status": "error"}), 400

        # Ensure syllabus_result is properly formatted as a list of strings
        if isinstance(syllabus_result, str):
            syllabus_result = [syllabus_result]
        elif not isinstance(syllabus_result, list):
            syllabus_result = [str(syllabus_result)]

        # Filter out any None values and ensure all items are strings
        syllabus_result = [str(item) for item in syllabus_result if item is not None]

        response_data = {
            "message": "PDF processed successfully.",
            "study_plan": plan if isinstance(plan, list) else [str(plan)],
            "syllabus": syllabus_result,
            "status": "success"
        }

        return jsonify(response_data), 200
    except Exception as e:
        logging.error(f"Error in upload_pdf route: {str(e)}")
        return jsonify({"message": f"Server error while processing PDF: {str(e)}", "status": "error"}), 500

@app.route('/generate_mcq', methods=['GET'])
@login_required
def get_mcq():
    """Generate MCQs from the stored document and return syllabus."""
    try:
        global current_document_id, syllabus, current_document_text
        if not current_document_id:
            return jsonify({"message": "Please upload a PDF first.", "status": "error"}), 400
        
        if not current_document_text:
            return jsonify({"message": "No document content available (Qdrant storage skipped).", "status": "error"}), 400
        
        questions = generate_mcq(current_document_text)
        
        # Ensure questions is a list
        if not isinstance(questions, list):
            questions = [questions]
        
        # Filter out any invalid questions
        valid_questions = []
        for q in questions:
            if isinstance(q, dict) and "question" in q and "options" in q and "correct" in q:
                valid_questions.append(q)
        
        response_data = {
            "message": f"Successfully generated {len(valid_questions)} MCQs." if valid_questions else "No valid MCQs generated.",
            "questions": valid_questions,
            "syllabus": syllabus if isinstance(syllabus, list) else [str(syllabus)] if syllabus else [],
            "status": "success" if valid_questions else "error"
        }
        
        return jsonify(response_data), 200 if valid_questions else 400
    except Exception as e:
        logging.error(f"Error in get_mcq route: {str(e)}")
        return jsonify({
            "message": f"Server error while generating MCQs: {str(e)}",
            "questions": [],
            "syllabus": [],
            "status": "error"
        }), 500

@app.route('/generate_mcq', methods=['POST'])
@login_required
def generate_mcq_route():
    """Generate MCQs (POST method)."""
    try:
        global current_document_text, current_document_id, mcq_questions
        if not current_document_text:
            logging.warning("MCQ generation attempted without document")
            return jsonify({"message": "No document available. Please upload a PDF first.", "status": "error"}), 400

        logging.info("Starting MCQ generation")
        questions = generate_mcq(current_document_text)
        logging.info(f"Generated {len(questions) if questions else 0} questions")
        
        # Validate questions format
        if not questions or not isinstance(questions, list):
            logging.error("No questions generated or invalid format")
            return jsonify({
                "message": "No suitable content found for generating MCQs.",
                "questions": [],
                "status": "error"
            }), 400
            
        # Filter out any invalid questions and keep only well-formed ones
        valid_questions = []
        for q in questions:
            if isinstance(q, dict) and all(key in q for key in ["question", "options", "correct"]):
                if q["options"] and isinstance(q["options"], list):
                    valid_questions.append(q)
                    
        if not valid_questions:
            logging.warning("No valid questions after filtering")
            return jsonify({
                "message": "Could not generate valid MCQs from the document content.",
                "questions": [],
                "status": "error"
            }), 400

        mcq_questions = valid_questions  # Update global questions list
        logging.info(f"Successfully validated and filtered {len(valid_questions)} MCQs")

        return jsonify({
            "message": f"Successfully generated {len(valid_questions)} MCQs.",
            "questions": valid_questions,
            "status": "success"
        }), 200
    except Exception as e:
        logging.error(f"Error in generate_mcq route: {str(e)}")
        return jsonify({
            "message": f"Server error while generating MCQs: {str(e)}",
            "questions": [],
            "status": "error"
        }), 500

@app.route('/generate_ppt', methods=['POST'])
@login_required
def generate_ppt():
    """Generate structured PPT slide content from uploaded PDF using GroQ."""
    try:
        global current_document_text, current_document_id, current_ppt_data
        if not current_document_text:
            return jsonify({
                "message": "No document available. Please upload a PDF first.",
                "status": "error"
            }), 400

        text = current_document_text[:5000]
        prompt = f"""
You are an expert educational content creator. A user has uploaded a PDF document to help generate a PowerPoint presentation. Your task is to create an in-depth, educational presentation.

Please analyze the document content below and generate **5 comprehensive slides**. For each slide:

1. Create a descriptive title that clearly indicates the topic
2. Generate 3-5 key points or concepts, where each point:
   - Highlights a significant concept or insight from the document
   - Uses clear, precise language
   - Is detailed enough to be educational (up to 150 characters per point)
   - Builds on previous points when possible

Guidelines:
- Focus on the most important concepts from the document
- Ensure points flow logically within each slide
- Use academic but accessible language
- Include specific examples or data points from the document
- Highlight relationships between concepts when relevant

Use this exact format:

Slide 1: [Clear, Descriptive Title]
- [Key Point 1: Include specific details or examples]
- [Key Point 2: Build on previous point if relevant]
- [Key Point 3: Add supporting information]
- [Key Point 4: Optional additional insight]
- [Key Point 5: Optional conclusion or next steps]

Slide 2: [Title that connects to overall flow]
- ...

Here is the uploaded PDF content:

{text}
"""
        completion = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at creating educational presentations."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=1500
        )
        ppt_text = completion.choices[0].message.content.strip()

        ppt_lines = [line.strip() for line in ppt_text.split("\n") if line.strip()]
        slides = []
        current_slide = {}

        for line in ppt_lines:
            if line.lower().startswith("slide"):
                if current_slide:
                    slides.append(current_slide)
                slide_title = line.split(":", 1)[-1].strip()
                current_slide = {"title": slide_title, "bullets": [], "details": []}
            elif line.startswith("-"):
                bullet = line[1:].strip()
                current_slide["bullets"].append(bullet)
                # Generate a detailed explanation for each bullet point
                try:
                    detail_prompt = f"Based on the uploaded document content, provide a detailed explanation (2-3 sentences) for this bullet point: '{bullet}'. Keep it student-friendly and focused on the topic. Context:\n{text[:200]}..."
                    detail_completion = groq_client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=[
                            {"role": "system", "content": "You are an expert at explaining educational concepts clearly."},
                            {"role": "user", "content": detail_prompt}
                        ],
                        temperature=0.7,
                        max_tokens=200
                    )
                    detail = detail_completion.choices[0].message.content.strip()
                    current_slide["details"].append(detail)
                except Exception as e:
                    logging.warning(f"Failed to generate detail for bullet point: {e}")
                    current_slide["details"].append("")

        if current_slide:
            slides.append(current_slide)

        # Create the PowerPoint presentation
        prs = Presentation()

        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Generated Presentation"
        subtitle.text = f"Created on {datetime.now().strftime('%B %d, %Y')}"

        # Add content slides
        for slide_data in slides:
            content_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = content_slide.shapes.title
            content = content_slide.placeholders[1]
            
            title.text = slide_data['title']
            
            # Format content with bullets and details
            content_text = []
            for bullet, detail in zip(slide_data['bullets'], slide_data['details']):
                content_text.append(f"• {bullet}")
                if detail:
                    content_text.append(f"  {detail}")
            
            content.text = "\n".join(content_text)

        # Save the presentation to BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        # Store the PPT data globally for download
        global current_ppt_data
        current_ppt_data = output.getvalue()

        return jsonify({
            "message": "Presentation generated successfully. Click Download PPT to save.",
            "status": "success"
        }), 200
    except Exception as e:
        logging.error(f"Error in generate_ppt route: {str(e)}")
        return jsonify({
            "message": f"Error generating PPT: {str(e)}",
            "status": "error"
        }), 500

@app.route('/download_ppt', methods=['POST'])
@login_required
def download_ppt():
    """Generate and download a PowerPoint presentation with predefined slides."""
    try:
        global current_document_id, current_ppt_data
        if not current_document_id:
            return jsonify({
                "message": "Please upload a PDF first.",
                "status": "error"
            }), 400

        if not current_ppt_data:
            return jsonify({
                "message": "Please generate the presentation first.",
                "status": "error"
            }), 400

        # Accept optional filename from client (JSON body). If not provided, use a default.
        filename = None
        try:
            data = request.get_json()
            filename = data.get('filename') if data else None
        except Exception:
            filename = None

        # Sanitize filename and ensure it ends with .pptx
        def sanitize_filename(name):
            # Replace invalid characters with underscore
            name = re.sub(r'[<>:"/\\|?*]', '_', name)
            # Remove any leading/trailing spaces or dots
            name = name.strip('. ')
            # Ensure the name isn't empty and has .pptx extension
            if not name:
                name = 'presentation'
            if not name.lower().endswith('.pptx'):
                name += '.pptx'
            return name

        download_name = sanitize_filename(filename or 'presentation.pptx')
        
        return send_file(
            BytesIO(current_ppt_data),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=download_name
        )
    except Exception as e:
        logging.error(f"Error in download_ppt route: {str(e)}")
        return jsonify({
            "message": f"Error downloading presentation: {str(e)}",
            "status": "error"
        }), 500


@app.route('/clarify_doubt', methods=['POST'])
@login_required
def clarify():
    """Clarify doubts using document content, syllabus, and Gemini."""
    try:
        global current_document_id
        if not current_document_id:
            return jsonify({"message": "Please upload a PDF first.", "status": "error"}), 400
        
        data = request.json
        if not data or 'question' not in data:
            return jsonify({"message": "Please provide a question.", "status": "error"}), 400
        
        question = data.get('question', '')
        response_type = data.get('response_type', 'short')
        if not question:
            return jsonify({"message": "Please provide a question.", "status": "error"}), 400
        
        response = clarify_doubts_with_rag(question, current_document_id, response_type)
        if isinstance(response, dict):
            return response
        return jsonify({"message": response, "status": "success"}), 200
    except Exception as e:
        logging.error(f"Error in clarify_doubt route: {str(e)}")
        return jsonify({"message": f"Server error while clarifying doubt: {str(e)}", "status": "error"}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)